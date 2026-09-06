"""
ingest_document.py
------------------
Document ingestion pipeline for the ORKA Equipment Knowledge Library.

Extracts text from uploaded PDF, Word, and PowerPoint files, splits it into
overlapping chunks, generates embeddings, and stores the chunks for retrieval.
The pipeline reads from the uploaded file object instead of copying the whole
upload into memory first, which is important on low-memory hosts.
"""

import logging
import os

import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation

from .local_db import test_db
from .vector_embed import embed_text

logger = logging.getLogger(__name__)

CHUNK_SIZE = 2000
CHUNK_OVERLAP = 200
MAX_EXTRACTED_CHARS = int(os.getenv("MAX_EXTRACTED_CHARS", "300000"))

SUPPORTED_TYPES = {".pdf", ".docx", ".pptx"}


def _seek_start(file_obj) -> None:
    file_obj.seek(0)


def extract_text_pdf(file_obj) -> str:
    """Extract text from a PDF. Returns empty string for scanned PDFs."""
    _seek_start(file_obj)
    text_parts = []
    total_chars = 0
    with pdfplumber.open(file_obj) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                remaining = MAX_EXTRACTED_CHARS - total_chars
                if remaining <= 0:
                    break
                page_text = page_text[:remaining]
                text_parts.append(page_text)
                total_chars += len(page_text)
    return "\n".join(text_parts)


def extract_text_docx(file_obj) -> str:
    """Extract text from a Word document paragraph by paragraph. hello"""
    _seek_start(file_obj)
    doc = DocxDocument(file_obj)
    text_parts = []
    total_chars = 0
    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        remaining = MAX_EXTRACTED_CHARS - total_chars
        if remaining <= 0:
            break
        text = text[:remaining]
        text_parts.append(text)
        total_chars += len(text)
    return "\n".join(text_parts)


def extract_text_pptx(file_obj) -> str:
    """Extract text from a PowerPoint file slide by slide."""
    _seek_start(file_obj)
    prs = Presentation(file_obj)
    text_parts = []
    total_chars = 0
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())
        if slide_texts:
            slide_text = f"[Slide {slide_num}]\n" + "\n".join(slide_texts)
            remaining = MAX_EXTRACTED_CHARS - total_chars
            if remaining <= 0:
                break
            slide_text = slide_text[:remaining]
            text_parts.append(slide_text)
            total_chars += len(slide_text)
    return "\n\n".join(text_parts)


def extract_text(filename: str, file_obj) -> str:
    """Route to the correct extractor based on file extension."""
    ext = filename.lower().rsplit(".", 1)[-1]
    if ext == "pdf":
        return extract_text_pdf(file_obj)
    if ext == "docx":
        return extract_text_docx(file_obj)
    if ext == "pptx":
        return extract_text_pptx(file_obj)
    raise ValueError(f"Unsupported file type: .{ext}")


def chunk_text(text: str):
    """Yield overlapping chunks without building a second list in memory."""
    start = 0
    while start < len(text):
        end = start + CHUNK_SIZE
        yield text[start:end]
        start += CHUNK_SIZE - CHUNK_OVERLAP


def ingest_document(filename: str, file_obj, equipment_tag: str, uploaded_by: str) -> dict:
    """
    Ingest one PDF, Word document, or PowerPoint file from a file-like object.

    The upload object is expected to be seekable, which FastAPI's UploadFile
    provides via a spooled temporary file.
    """
    ext = "." + filename.lower().rsplit(".", 1)[-1]
    if ext not in SUPPORTED_TYPES:
        return {
            "filename": filename,
            "status": "skipped",
            "reason": f"Unsupported file type: {ext}",
            "chunks_stored": 0,
            "chunks_failed": 0,
        }

    text = extract_text(filename, file_obj)
    extracted_chars = len(text)
    truncated = extracted_chars >= MAX_EXTRACTED_CHARS
    if not text.strip():
        logger.warning("No extractable text found in %s.", filename)
        return {
            "filename": filename,
            "status": "skipped",
            "reason": "No extractable text - scanned or image-only file?",
            "chunks_stored": 0,
            "chunks_failed": 0,
            "extracted_chars": 0,
            "truncated": False,
        }

    stored = 0
    failed = 0
    for i, chunk in enumerate(chunk_text(text)):
        try:
            embedding = embed_text(chunk)
            test_db.insert_document_chunk(
                filename=filename,
                file_type=ext.lstrip("."),
                content_chunk=chunk,
                chunk_index=i,
                equipment_tag=equipment_tag,
                embedding=embedding,
                uploaded_by=uploaded_by,
            )
            stored += 1
        except Exception as e:
            logger.warning("Failed to store chunk %d of %s: %s", i, filename, e)
            failed += 1

    logger.info("Ingested %s: %d chunks stored, %d failed.", filename, stored, failed)
    return {
        "filename": filename,
        "status": "complete",
        "chunks_stored": stored,
        "chunks_failed": failed,
        "extracted_chars": extracted_chars,
        "truncated": truncated,
    }
